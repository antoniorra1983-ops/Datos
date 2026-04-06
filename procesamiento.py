import pandas as pd
import numpy as np
import re
import holidays
from io import BytesIO
from datetime import datetime, date, timedelta, time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

chile_holidays = holidays.Chile()

DISTANCIAS = {
    "PU-LI": 43.13, "LI-PU": 43.13, "PU-SA": 29.11, "SA-PU": 29.11,
    "EB-PU": 25.40, "PU-EB": 25.40, "VM-LI": 34.03, "LI-VM": 34.03,
    "VM-PU": 9.10,  "PU-VM": 9.10
}

def parse_latam_number(val):
    if pd.isna(val): return 0.0
    if isinstance(val, (int, float)): return float(val)
    s = str(val).strip().replace(' ', '').replace('$', '')
    s = re.sub(r'[^\d.,-]', '', s)
    if not s: return 0.0
    if ',' in s and '.' in s:
        if s.rfind(',') > s.rfind('.'): s = s.replace('.', '').replace(',', '.')
        else: s = s.replace(',', '')
    elif ',' in s: s = s.replace(',', '.')
    try: return float(s)
    except: return 0.0

def get_tipo_dia(fch):
    if fch in chile_holidays or fch.weekday() == 6: return "D/F"
    if fch.weekday() == 5: return "S"
    return "L"

def convertir_a_minutos(val):
    if pd.isna(val) or str(val).strip() == "": return None
    try:
        if isinstance(val, (datetime, time)): return val.hour * 60 + val.minute + (val.second / 60.0)
        if isinstance(val, str):
            val = val.strip()
            m_ss = re.search(r'(\d{1,2}):(\d{2}):(\d{2})', val)
            if m_ss: return int(m_ss.group(1)) * 60 + int(m_ss.group(2)) + (int(m_ss.group(3)) / 60.0)
            m_mm = re.search(r'(\d{1,2}):(\d{2})', val)
            if m_mm: return int(m_mm.group(1)) * 60 + int(m_mm.group(2))
        return None
    except: return None

def format_hms(minutos_float, con_signo=False):
    if pd.isna(minutos_float) or minutos_float == 0: return "00:00:00"
    signo = ("+" if minutos_float > 0 else "-" if minutos_float < 0 else "") if con_signo else ""
    total_segundos = int(round(abs(minutos_float) * 60))
    h, r = divmod(total_segundos, 3600)
    m, s = divmod(r, 60)
    return f"{signo}{h:02d}:{m:02d}:{s:02d}"

def clasificar_flota_func(motriz):
    try:
        num = int(float(motriz))
        if 1 <= num <= 27: return "XT-100"
        if 28 <= num <= 35: return "XT-M"
        if 101 <= num <= 110: return "SFE (Chino)"
        return "OTRO"
    except: return "S/I"

def leer_fecha_archivo(file):
    try:
        df = pd.read_excel(file, nrows=1, header=None)
        val = str(df.iloc[0, 0]).split('.')[0].strip().zfill(6)
        if len(val) == 6 and val.isdigit():
            return datetime(2000 + int(val[4:6]), int(val[2:4]), int(val[0:2])).date()
        match = re.search(r'(\d{2})[.-](\d{2})[.-](\d{2})', file.name)
        if match: return datetime(2000 + int(match.group(3)), int(match.group(2)), int(match.group(1))).date()
    except: pass
    return None

def procesar_thdr_avanzado(file):
    try:
        df_raw = pd.read_excel(file, header=None)
        est_h = df_raw.iloc[0].ffill().values
        df = df_raw.iloc[2:].copy()
        df = df[df.iloc[:, 1].notna()] # Más flexible para el número de servicio
        columnas_base = ["Recorrido", "Servicio", "Hora_Prog", "Motriz 1", "Motriz 2", "Unidad"]
        estaciones_raw = [str(est_h[i]) if pd.notna(est_h[i]) else f"Col_{i}" for i in range(6, len(df_raw.columns))]
        nombres_finales, conteos = list(columnas_base), {}
        for nombre in estaciones_raw:
            if nombre not in conteos:
                conteos[nombre] = 0; nombres_finales.append(nombre)
            else:
                conteos[nombre] += 1; nombres_finales.append(f"{nombre}.{conteos[nombre]}")
        df.columns = nombres_finales[:len(df.columns)]
        est_cols = df.columns[6:]
        def get_trip(row):
            h_reales = row[est_cols].apply(convertir_a_minutos).dropna()
            if len(h_reales) < 2: return "OTRO", "OTRO", None, 0, 0
            def cod(n_est):
                n = str(n_est).upper()
                if "PUERTO" in n: return "PU"
                if "LIMACHE" in n: return "LI"
                if "SA" in n or "ALDEA" in n: return "SA"
                return n[:2]
            t_s, t_l = h_reales.iloc[0], h_reales.iloc[-1]
            if t_l < t_s: t_l += 1440
            return (f"{cod(h_reales.index[0])}-{cod(h_reales.index[-1])}", cod(h_reales.index[0]), t_s, (t_l - t_s), int(t_s // 60) % 24)
        stats = df.apply(get_trip, axis=1)
        df['Tipo_Rec'], df['Origen'], df['Min_S_Real'], df['TDV_Min'], df['Hora_Salida'] = ([x[0] for x in stats], [x[1] for x in stats], [x[2] for x in stats], [x[3] for x in stats], [x[4] for x in stats])
        df['Min_Prog'] = df['Hora_Prog'].apply(convertir_a_minutos)
        df['Retraso'] = df['Min_S_Real'] - df['Min_Prog']
        df['Puntual'] = (abs(df['Retraso']) <= 5).astype(int)
        df['Dist_Base'] = df['Tipo_Rec'].map(DISTANCIAS).fillna(0)
        df['Peso'] = df['Unidad'].apply(lambda x: 2 if str(x).strip().upper() == 'M' else 1)
        df['Tren-Km'] = df['Dist_Base'] * df['Peso']
        df['Flota'] = df['Motriz 1'].apply(clasificar_flota_func)
        f_dt = leer_fecha_archivo(file)
        df['Fecha_Op'] = f_dt if f_dt else None
        return df, df['Tren-Km'].sum(), df[df['TDV_Min'] > 0]['TDV_Min'].mean(), (df['Puntual'].sum() / len(df) * 100) if len(df) > 0 else 0
    except Exception as e: return pd.DataFrame(), 0, 0, 0

def to_pptx(title_text, df=None, metrics_dict=None):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"EFE Valparaíso: {title_text}"
    y_cursor = Inches(1.5)
    if metrics_dict:
        tf = slide.shapes.add_textbox(Inches(0.5), y_cursor, Inches(9), Inches(1)).text_frame
        for k, v in metrics_dict.items():
            p = tf.add_paragraph()
            p.text, p.font.size, p.font.bold = f"• {k}: {v}", Pt(16), True
            p.font.color.rgb = RGBColor(0, 81, 149)
        y_cursor += Inches(1.2)
    if df is not None and not df.empty:
        df_display = df.head(12).reset_index(drop=True)
        rows, cols = df_display.shape
        table = slide.shapes.add_table(rows+1, cols, Inches(0.5), y_cursor, Inches(9), Inches(3)).table
        for c, col_name in enumerate(df_display.columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 81, 149)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        for r in range(rows):
            for c in range(cols):
                val = df_display.iloc[r, c]
                table.cell(r+1, c).text = str(val) if not isinstance(val, float) else f"{val:,.1f}"
    output = BytesIO()
    prs.save(output)
    return output.getvalue()

def to_excel_consolidado(df_ops, df_tr, df_tr_acum, df_seat, df_p_d, df_p_15, df_fact_h, df_fact_d):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        dfs = {'Operaciones': df_ops, 'Kms_Diarios_Tren': df_tr, 'Odometros_Acum_Tren': df_tr_acum, 'SEAT': df_seat, 'PRMTE_D': df_p_d, 'PRMTE_15': df_p_15, 'Fact_H': df_fact_h, 'Fact_D': df_fact_d}
        for name, df in dfs.items():
            if not df.empty: df.to_excel(writer, index=False, sheet_name=name)
    return output.getvalue()

def exportar_resumen_excel(metrics_dict, df_resumen_jornada, df_energia):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        pd.DataFrame([metrics_dict]).T.to_excel(writer, sheet_name='Métricas')
        if not df_resumen_jornada.empty: df_resumen_jornada.to_excel(writer, sheet_name='Resumen_Jornada')
        if not df_energia.empty: df_energia.to_excel(writer, sheet_name='Energía')
    return output.getvalue()
